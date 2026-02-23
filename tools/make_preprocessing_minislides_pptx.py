from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


def add_bullets(slide, title: str, bullets: list[str]) -> None:
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    first = True
    for bullet in bullets:
        if first:
            body.text = bullet
            first = False
            continue
        p = body.add_paragraph()
        p.text = bullet
        p.level = 0


def add_notes(slide, notes: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = notes.strip()


def build_presentation(out_path: Path) -> None:
    prs = Presentation()

    # Slide 1: CV preprocessing + detection
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    add_bullets(
        slide1,
        title="Preprocessing for Price-Box Detection (OpenCV)",
        bullets=[
            "Convert to grayscale + CLAHE (local contrast)",
            "Adaptive thresholding → binary text map",
            "Morphology to connect characters → rectangular blocks",
            "Color masks (HSV + LAB) + connected components → candidate boxes",
            "Heuristic filtering + IoU-merge → final price boxes",
        ],
    )

    add_notes(
        slide1,
        notes=(
            "Goal: robustly find the flyer’s price-label rectangles before OCR.\n\n"
            "1) Text-oriented preprocessing (src/ocr/pipeline/image_preprocessing.py):\n"
            "- BGR → grayscale (to_grayscale).\n"
            "- CLAHE (enhance_contrast): clipLimit=2.0, tile grid 8×8 to boost local contrast without blowing up noise.\n"
            "- Adaptive Gaussian threshold (adaptive_binarize): block_size=21, C=8 to handle uneven lighting/background.\n\n"
            "2) Two different morphological outputs are produced:\n"
            "- morph_text (morph_refine): close (3×3) then open (2×2) to clean small gaps/noise for OCR.\n"
            "- morph (morph_for_pricebox_detection): designed for detection, not OCR.\n"
            "  * Invert binary so text becomes white.\n"
            "  * Dilate horizontally with kernel (9×2) to connect characters into words/prices.\n"
            "  * Dilate vertically with kernel (2×5) to connect multi-line labels into blocks.\n"
            "  * Close with kernel (5×5) to fill small holes.\n\n"
            "3) Color-based candidate generation (preprocess_for_detection / pricebox_detection):\n"
            "- HSV mask for white boxes + LAB L-channel threshold (bright regions).\n"
            "- HSV masks for blue regions and dark (text) regions; combine into a text_mask.\n"
            "- Connected components (cv2.connectedComponentsWithStats) to get candidate rectangles.\n"
            "- Heuristic filtering (area, rectangularity, aspect ratio, etc.).\n"
            "- Merge overlaps using IoU clustering (IOU_MERGE_THRESHOLD = 0.3).\n\n"
            "Result: a small set of bounding boxes (x,y,w,h) that should contain the price tag graphics."
        ),
    )

    # Slide 2: OCR + NLP extraction
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    add_bullets(
        slide2,
        title="From Price Box → OCR Text → NLP Product Name",
        bullets=[
            "Crop multiple regions around each detected price box",
            "Run Tesseract with several PSM modes + light binarization",
            "Keep ALL readable lines (not just one) as OCR text",
            "spaCy (de_core_news_sm) selects the likely product name from OCR text",
            "Output includes both full OCR text and extracted product name",
        ],
    )

    add_notes(
        slide2,
        notes=(
            "Goal: make the pipeline ‘NLP-based’ by separating (a) getting noisy text via OCR from (b) deciding what is the product name via NLP.\n\n"
            "1) Region selection around the price box (src/ocr/pipeline/ocr_product_text.py):\n"
            "- The detected box mostly contains the price, not the product title.\n"
            "- Therefore multiple nearby crops are tested (e.g., left / above / left-above / right), and the best OCR candidate is selected by a scoring heuristic.\n\n"
            "2) OCR strategy (Tesseract via pytesseract):\n"
            "- Upscale crop (scale ≈ 3×) with cubic interpolation to help Tesseract on small flyer fonts.\n"
            "- Convert to grayscale; try Otsu threshold and also an inverted variant (to handle dark/white-on-color cases).\n"
            "- Run multiple Tesseract page segmentation modes: psm 7 (single line), 6 (block), 4 (column), 11 (sparse).\n"
            "- Instead of picking one ‘best line’, keep all reasonably clean lines and concatenate them → this becomes the OCR text context.\n"
            "  (Reason: the true product name might be one line among specs / units / marketing text.)\n\n"
            "3) NLP product extraction (src/extraction/ner_model.py):\n"
            "- spaCy German model de_core_news_sm provides tokenization + POS/NER signals.\n"
            "- Heuristics bias toward: known brands, noun-rich phrases, title-like capitalization, and reject disclaimer-only text.\n"
            "- If NLP rejects, product_name stays empty (no fallback to ‘full OCR text as product name’).\n\n"
            "Result: for each price box we store\n"
            "- full OCR text (context)\n"
            "- extracted product name (NLP decision)\n"
            "- price (separate OCR specialized for digits)"
        ),
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Create a minimal 1–2 slide PPTX with detailed presenter notes.")
    parser.add_argument(
        "--out",
        type=Path,
        default=Path("preprocessing_minislides.pptx"),
        help="Output PPTX path (default: preprocessing_minislides.pptx)",
    )
    args = parser.parse_args()

    build_presentation(args.out)
    print(f"Wrote: {args.out.resolve()}")


if __name__ == "__main__":
    main()
