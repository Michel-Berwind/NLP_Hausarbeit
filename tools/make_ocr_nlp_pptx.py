from __future__ import annotations

import json
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple

import cv2
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WORKSPACE_ROOT = Path(__file__).resolve().parents[1]


def _safe_add_picture(slide, image_path: Path, left, top, width=None, height=None) -> bool:
	if not image_path.exists():
		return False
	try:
		if width is not None and height is not None:
			slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
		elif width is not None:
			slide.shapes.add_picture(str(image_path), left, top, width=width)
		elif height is not None:
			slide.shapes.add_picture(str(image_path), left, top, height=height)
		else:
			slide.shapes.add_picture(str(image_path), left, top)
		return True
	except Exception:
		return False


def _add_caption(slide, text: str, left, top, width, height, font_size_pt: int = 12) -> None:
	box = slide.shapes.add_textbox(left, top, width, height)
	p = box.text_frame.paragraphs[0]
	p.text = text
	p.font.size = Pt(font_size_pt)
	p.alignment = PP_ALIGN.CENTER


def _add_bullets(slide, title: str, bullets: Sequence[str], note: Optional[str] = None) -> None:
	if slide.shapes.title:
		slide.shapes.title.text = title
	body = None
	for shape in slide.shapes:
		if shape.has_text_frame and shape != slide.shapes.title:
			body = shape
			break
	if body is None:
		body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(5.5))

	frame = body.text_frame
	frame.clear()
	for i, b in enumerate(bullets):
		p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
		p.text = b
		p.level = 0
		p.font.size = Pt(20)
	if note:
		p = frame.add_paragraph()
		p.text = note
		p.level = 0
		p.font.size = Pt(14)


def _read_text_excerpt(path: Path, max_lines: int = 12) -> str:
	if not path.exists():
		return ""
	lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
	lines = [ln.rstrip() for ln in lines if ln.strip()]
	return "\n".join(lines[:max_lines])


def _compute_product_regions(box: List[int], img_shape: Tuple[int, int, int]) -> Dict[str, Tuple[int, int, int, int]]:
	"""Replicates the region geometry from extract_product_text (ocr_product_text.py)."""
	x, y, w, h = box
	H, W = img_shape[:2]
	regions: Dict[str, Tuple[int, int, int, int]] = {}

	# ABOVE
	x0_above = max(0, x - int(0.5 * w))
	x1_above = min(W, x + w + int(0.5 * w))
	y1_above = max(0, y - int(0.1 * h))
	y0_above = max(0, y - int(2.0 * h))
	if y1_above - y0_above >= 20 and x1_above - x0_above >= 40:
		regions["above"] = (x0_above, y0_above, x1_above, y1_above)

	# LEFT
	x1_left = max(0, x - int(0.1 * w))
	x0_left = max(0, x - int(2.0 * w))
	y0_left = max(0, y - int(0.5 * h))
	y1_left = min(H, y + h + int(0.5 * h))
	if x1_left - x0_left >= 40 and y1_left - y0_left >= 20:
		regions["left"] = (x0_left, y0_left, x1_left, y1_left)

	# RIGHT
	x0_right = min(W, x + w + int(0.1 * w))
	x1_right = min(W, x + w + int(2.0 * w))
	y0_right = max(0, y - int(0.5 * h))
	y1_right = min(H, y + h + int(0.5 * h))
	if x1_right - x0_right >= 40 and y1_right - y0_right >= 20:
		regions["right"] = (x0_right, y0_right, x1_right, y1_right)

	# BELOW
	x0_below = max(0, x - int(0.5 * w))
	x1_below = min(W, x + w + int(0.5 * w))
	y0_below = min(H, y + h + int(0.1 * h))
	y1_below = min(H, y + h + int(1.0 * h))
	if y1_below - y0_below >= 20 and x1_below - x0_below >= 40:
		regions["below"] = (x0_below, y0_below, x1_below, y1_below)

	return regions


def _draw_regions_overlay(img_bgr, box: List[int], regions: Dict[str, Tuple[int, int, int, int]]):
	"""Returns an overlay image with the price box + product search regions drawn."""
	x, y, w, h = box
	out = img_bgr.copy()
	# Price box
	cv2.rectangle(out, (x, y), (x + w, y + h), (0, 255, 0), 6)
	cv2.putText(out, "price box", (x, max(30, y - 10)), cv2.FONT_HERSHEY_SIMPLEX, 1.2, (0, 255, 0), 3)

	colors = {
		"above": (255, 0, 0),
		"left": (0, 128, 255),
		"right": (255, 128, 0),
		"below": (255, 0, 255),
	}
	for name, (x0, y0, x1, y1) in regions.items():
		c = colors.get(name, (200, 200, 200))
		cv2.rectangle(out, (x0, y0), (x1, y1), c, 4)
		cv2.putText(out, name, (x0 + 10, y0 + 40), cv2.FONT_HERSHEY_SIMPLEX, 1.2, c, 3)

	return out


def _save_region_crops(img_bgr, regions: Dict[str, Tuple[int, int, int, int]], out_dir: Path) -> Dict[str, Path]:
	out_paths: Dict[str, Path] = {}
	for name, (x0, y0, x1, y1) in regions.items():
		crop = img_bgr[y0:y1, x0:x1]
		p = out_dir / f"region_{name}.png"
		cv2.imwrite(str(p), crop)
		out_paths[name] = p
	return out_paths


def _json_snippet(path: Path, max_chars: int = 900) -> str:
	if not path.exists():
		return ""
	text = path.read_text(encoding="utf-8", errors="ignore")
	text = "\n".join(text.splitlines()[:60])
	return text[:max_chars] + ("…" if len(text) > max_chars else "")


def build_presentation(output_path: Path) -> None:
	prs = Presentation()

	img_page = WORKSPACE_ROOT / "data" / "images" / "aldi" / "page_coffee.png"
	ocr_page_txt = WORKSPACE_ROOT / "data" / "ocr_text" / "aldi" / "page_coffee.txt"
	anno_json = WORKSPACE_ROOT / "data" / "annotations" / "page_coffee.json"

	debug_out = WORKSPACE_ROOT / "data" / "debug_ocr_slides" / "page_coffee"
	debug_out.mkdir(parents=True, exist_ok=True)

	# Load one example detection (first item)
	box: Optional[List[int]] = None
	product: Optional[str] = None
	price: Optional[str] = None
	if anno_json.exists():
		data = json.loads(anno_json.read_text(encoding="utf-8"))
		if data and isinstance(data, list) and data[0].get("items"):
			first = data[0]["items"][0]
			box = first.get("box")
			product = first.get("product")
			price = first.get("price")

	# Create overlay + region crops
	overlay_path = debug_out / "product_text_regions_overlay.png"
	region_paths: Dict[str, Path] = {}
	if img_page.exists() and box:
		img_bgr = cv2.imread(str(img_page))
		regions = _compute_product_regions(box, img_bgr.shape)
		overlay = _draw_regions_overlay(img_bgr, box, regions)
		cv2.imwrite(str(overlay_path), overlay)
		region_paths = _save_region_crops(img_bgr, regions, debug_out)

	# Example crop paths
	crop_dir_blue = WORKSPACE_ROOT / "data" / "debug_price_boxes" / "page_coffee" / "crops"
	crop_dir_white = WORKSPACE_ROOT / "data" / "debug_price_boxes" / "page_coffee" / "white_boxes" / "crops"
	crop_examples = []
	for p in [
		crop_dir_blue / "box_00.png",
		crop_dir_blue / "box_01.png",
		crop_dir_white / "box_00.png",
		crop_dir_white / "box_01.png",
	]:
		if p.exists():
			crop_examples.append(p)

	# --- Slide 1: Title ---
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes.title.text = "OCR & NLP (post-detection pipeline)"
	sub = slide.placeholders[1]
	sub.text = "Example: page_coffee"
	_safe_add_picture(slide, img_page, Inches(8.0), Inches(1.7), width=Inches(5.0))

	# --- Slide 2: OCR overview ---
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	_add_bullets(
		slide,
		title="OCR tasks (two OCR problems)",
		bullets=[
			"Price OCR: read numeric price from each detected pricebox crop",
			"Product OCR: read product name text near each detected pricebox",
			"Uses multiple preprocessing attempts + multiple Tesseract PSM modes",
			"Outputs: (box, price, product) per detected item",
		],
		note="Code: src/ocr/pipeline/pricebox_detection.py + src/ocr/pipeline/ocr_product_text.py",
	)

	# --- Slide 3: Full page OCR (baseline) ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	if slide.shapes.title:
		slide.shapes.title.text = "Full-page OCR (baseline)"

	_safe_add_picture(slide, img_page, Inches(0.6), Inches(1.5), width=Inches(6.6))
	_add_caption(slide, "Input page image", Inches(0.6), Inches(5.8), Inches(6.6), Inches(0.3))

	excerpt = _read_text_excerpt(ocr_page_txt, max_lines=14)
	box_text = slide.shapes.add_textbox(Inches(7.4), Inches(1.5), Inches(5.7), Inches(4.6))
	frame = box_text.text_frame
	frame.word_wrap = True
	p = frame.paragraphs[0]
	p.text = "Example output (page_coffee.txt):\n\n" + (excerpt if excerpt else "(missing OCR text file)")
	p.font.size = Pt(14)

	_add_caption(slide, "Generated by src/ocr/ocr_single_page.py (Tesseract PSM 6)", Inches(7.4), Inches(6.2), Inches(5.7), Inches(0.4), font_size_pt=12)

	# --- Slide 4: Price OCR from crops ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	if slide.shapes.title:
		slide.shapes.title.text = "Price OCR from detected crops"

	# place up to 4 crops
	positions = [
		(Inches(0.7), Inches(1.5)),
		(Inches(4.0), Inches(1.5)),
		(Inches(7.3), Inches(1.5)),
		(Inches(10.6), Inches(1.5)),
	]
	for i, pth in enumerate(crop_examples[:4]):
		left, top = positions[i]
		_safe_add_picture(slide, pth, left, top, width=Inches(3.0))
		_add_caption(slide, pth.name, left, top + Inches(2.6), Inches(3.0), Inches(0.3), font_size_pt=11)

	bul = slide.shapes.add_textbox(Inches(0.7), Inches(4.4), Inches(13.0), Inches(3.6))
	fr = bul.text_frame
	fr.clear()
	p = fr.paragraphs[0]
	p.text = "How price OCR works (extract_price_from_crop):"
	p.font.size = Pt(18)
	for line in [
		"Resize crop ×4 and convert to grayscale",
		"Try several binarizations: inverted + Otsu, inverted + manual threshold, binary-inv Otsu",
		"Run Tesseract with multiple PSM modes (6, 7, 11, 13)",
		"Parse patterns like d.dd; validate range (0.10–99.99); small corrections (e.g., X.29 → X.99)",
	]:
		pp = fr.add_paragraph()
		pp.text = line
		pp.level = 1
		pp.font.size = Pt(16)

	# --- Slide 5: Product text OCR regions ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	if slide.shapes.title:
		slide.shapes.title.text = "Product text OCR: search regions around the price box"

	if overlay_path.exists():
		_safe_add_picture(slide, overlay_path, Inches(0.6), Inches(1.4), width=Inches(8.0))
		_add_caption(slide, "Green = price box; colored rectangles = candidate product-text regions", Inches(0.6), Inches(6.9), Inches(8.0), Inches(0.35), font_size_pt=12)

	# region crops on the right
	order = ["above", "left", "right", "below"]
	rpos = [
		(Inches(8.9), Inches(1.6)),
		(Inches(8.9), Inches(3.6)),
		(Inches(8.9), Inches(5.6)),
		(Inches(8.9), Inches(7.6)),
	]
	for name, (left, top) in zip(order, rpos):
		pth = region_paths.get(name)
		if pth and pth.exists():
			_safe_add_picture(slide, pth, left, top, width=Inches(4.2))
			_add_caption(slide, f"{name} region", left, top + Inches(1.65), Inches(4.2), Inches(0.25), font_size_pt=11)

	# --- Slide 6: Product OCR scoring & selection ---
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	_add_bullets(
		slide,
		title="Product OCR: best-text selection",
		bullets=[
			"For each region, run multiple preprocessing candidates (Otsu, adaptive, inverted)",
			"Run Tesseract with PSM 6/4/3 (block, column, auto)",
			"Clean lines + discard price-like lines (regex \\\"\\d+[.,]\\d{2}\\\")",
			"Score lines with heuristics (letters ratio, digits ratio, length, word structure)",
			"Pick the highest-scoring line across all regions",
		],
		note="Code: src/ocr/pipeline/ocr_product_text.py (score_line, extract_product_text)",
	)

	# --- Slide 7: NLP / structured output ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	if slide.shapes.title:
		slide.shapes.title.text = "NLP output: structured JSON for downstream processing"

	json_txt = _json_snippet(anno_json)
	box_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12.9), Inches(6.6))
	frame = box_text.text_frame
	frame.word_wrap = True
	p = frame.paragraphs[0]
	p.text = (
		"Pipeline output example (page_coffee.json):\n\n" + (json_txt if json_txt else "(missing JSON output)")
	)
	p.font.size = Pt(12)

	cap = "Code: src/ocr/pipeline/nlp_to_json.py (to_result_record + save_results)"
	if box and price is not None and product is not None:
		cap += f" | Example item: price={price}, product={product!r}"
	_add_caption(slide, cap, Inches(0.7), Inches(8.2), Inches(12.9), Inches(0.4), font_size_pt=11)

	output_path.parent.mkdir(parents=True, exist_ok=True)
	prs.save(str(output_path))


def main() -> None:
	out = WORKSPACE_ROOT / "ocr_and_nlp_pipeline.pptx"
	build_presentation(out)
	print(f"Wrote: {out}")


if __name__ == "__main__":
	main()
