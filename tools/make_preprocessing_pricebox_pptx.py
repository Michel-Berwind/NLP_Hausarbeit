from __future__ import annotations

from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WORKSPACE_ROOT = Path(__file__).resolve().parents[1]


def _safe_add_picture(slide, image_path: Path, left, top, width=None, height=None) -> bool:
	if not image_path.exists():
		return False
	try:
		if width is not None and height is not None:
			slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
		elif width is not None:
			slide.shapes.add_picture(str(image_path), left, top, width=width)
		elif height is not None:
			slide.shapes.add_picture(str(image_path), left, top, height=height)
		else:
			slide.shapes.add_picture(str(image_path), left, top)
		return True
	except Exception:
		return False


def _add_caption(slide, text: str, left, top, width, height, font_size_pt: int = 12) -> None:
	box = slide.shapes.add_textbox(left, top, width, height)
	p = box.text_frame.paragraphs[0]
	p.text = text
	p.font.size = Pt(font_size_pt)
	p.alignment = PP_ALIGN.CENTER


def _add_bullets(slide, title: str, bullets: Sequence[str], note: Optional[str] = None) -> None:
	layout = slide.slide_layout
	# Expecting a title+content layout; caller should create slide with layout 1.
	if slide.shapes.title:
		slide.shapes.title.text = title
	body = None
	for shape in slide.shapes:
		if shape.has_text_frame and shape != slide.shapes.title:
			body = shape
			break
	if body is None:
		# Fallback: add our own textbox
		body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(5.5))

	frame = body.text_frame
	frame.clear()
	for i, b in enumerate(bullets):
		p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
		p.text = b
		p.level = 0
		p.font.size = Pt(20)
	if note:
		p = frame.add_paragraph()
		p.text = note
		p.level = 0
		p.font.size = Pt(14)


def _add_image_grid(
	slide,
	title: str,
	images_and_labels: Sequence[Tuple[Path, str]],
	rows: int,
	cols: int,
	left: float = 0.6,
	top: float = 1.2,
	cell_w: float = 4.2,
	cell_h: float = 2.6,
	label_h: float = 0.35,
) -> None:
	if slide.shapes.title:
		slide.shapes.title.text = title

	max_items = rows * cols
	items = list(images_and_labels)[:max_items]

	for idx, (img_path, label) in enumerate(items):
		r = idx // cols
		c = idx % cols
		x = Inches(left + c * cell_w)
		y = Inches(top + r * cell_h)

		# Image area (reserve some space for label)
		img_h = Inches(cell_h - label_h)
		img_w = Inches(cell_w)
		added = _safe_add_picture(slide, img_path, x, y, width=img_w, height=img_h)
		if not added:
			# If missing, still show placeholder label
			pass

		_add_caption(
			slide,
			label + (" (missing)" if not img_path.exists() else ""),
			x,
			y + img_h,
			Inches(cell_w),
			Inches(label_h),
			font_size_pt=12,
		)


def build_presentation(output_path: Path) -> None:
	prs = Presentation()

	# Source assets (page_coffee example)
	img_page = WORKSPACE_ROOT / "data" / "images" / "aldi" / "page_coffee.png"
	pre_dir = WORKSPACE_ROOT / "data" / "debug_price_boxes" / "preprocess" / "page_coffee"
	det_dir = WORKSPACE_ROOT / "data" / "debug_price_boxes" / "page_coffee"
	white_det_dir = det_dir / "white_boxes"

	# --- Slide 1: Title ---
	slide = prs.slides.add_slide(prs.slide_layouts[0])
	slide.shapes.title.text = "Image Preprocessing & Pricebox Detection"
	sub = slide.placeholders[1]
	sub.text = "Example: page_coffee (ALDI flyer)"
	_safe_add_picture(slide, img_page, Inches(8.0), Inches(1.7), width=Inches(5.0))

	# --- Slide 2: Pipeline overview ---
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	_add_bullets(
		slide,
		title="End-to-end pipeline (what happens on one page)",
		bullets=[
			"Load page image (BGR)",
			"Preprocess: text-focused + detection masks (saved for debugging)",
			"Detect price boxes (blue boxes + white boxes)",
			"Crop each box and OCR the price",
			"(Optional) OCR product description above the price box",
		],
		note="Code: src/ocr/pipeline/pipeline_runner.py → image_preprocessing.py + pricebox_detection.py",
	)

	# --- Slide 3: Text-focused preprocessing ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
	_add_image_grid(
		slide,
		title="Image preprocessing (text-focused stages)",
		images_and_labels=[
			(pre_dir / "gray.png", "gray (grayscale)"),
			(pre_dir / "enhanced.png", "enhanced (CLAHE)"),
			(pre_dir / "binary.png", "binary (adaptive threshold)"),
			(pre_dir / "morph_text.png", "morph_text (close+open for OCR)"),
			(pre_dir / "morph.png", "morph (dilation/close for blocks)"),
		],
		rows=2,
		cols=3,
		cell_w=4.2,
		cell_h=2.8,
	)

	# --- Slide 4: Detection masks preprocessing ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	_add_image_grid(
		slide,
		title="Precomputed masks (detection-focused)",
		images_and_labels=[
			(pre_dir / "white_mask.png", "white_mask (bright background)"),
			(pre_dir / "blue_mask.png", "blue_mask (blue text)"),
			(pre_dir / "dark_mask.png", "dark_mask (dark text)"),
			(pre_dir / "text_mask.png", "text_mask = blue ∪ dark"),
			(pre_dir / "hsv.png", "HSV colorspace"),
			(pre_dir / "L.png", "LAB L channel (lightness)"),
		],
		rows=2,
		cols=3,
		cell_w=4.2,
		cell_h=2.8,
	)

	# --- Slide 5: Detection approach ---
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	_add_bullets(
		slide,
		title="Pricebox detection (how boxes are found)",
		bullets=[
			"Two-pass detection: (1) blue boxes with light text, (2) white boxes with colored/dark text",
			"Create box mask in HSV (+ LAB L threshold), then connected components",
			"Each component is a candidate box region",
			"Filter candidates with deterministic heuristics (size, aspect, text, brightness, contrast)",
			"Merge overlaps (IoU clustering), then sort in reading order",
		],
		note="Code: src/ocr/pipeline/pricebox_detection.py (detect_price_boxes, filter_white_box)",
	)

	# --- Slide 6: Masks and detection results (blue + white pass) ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	if slide.shapes.title:
		slide.shapes.title.text = "Detection visuals (page_coffee)"

	# Left: blue pass; Right: white pass
	_safe_add_picture(slide, det_dir / "01_box_mask.png", Inches(0.6), Inches(1.4), width=Inches(6.3))
	_add_caption(slide, "Blue-box pass: 01_box_mask.png", Inches(0.6), Inches(4.95), Inches(6.3), Inches(0.3))

	_safe_add_picture(slide, det_dir / "01_text_mask.png", Inches(0.6), Inches(5.3), width=Inches(6.3))
	_add_caption(slide, "Blue-box pass: 01_text_mask.png", Inches(0.6), Inches(8.85), Inches(6.3), Inches(0.3))

	_safe_add_picture(slide, white_det_dir / "01_box_mask_white.png", Inches(7.1), Inches(1.4), width=Inches(6.3))
	_add_caption(slide, "White-box pass: 01_box_mask_white.png", Inches(7.1), Inches(4.95), Inches(6.3), Inches(0.3))

	_safe_add_picture(slide, white_det_dir / "01_text_mask_white.png", Inches(7.1), Inches(5.3), width=Inches(6.3))
	_add_caption(slide, "White-box pass: 01_text_mask_white.png", Inches(7.1), Inches(8.85), Inches(6.3), Inches(0.3))

	# --- Slide 7: Filtering + merging (before/after) ---
	slide = prs.slides.add_slide(prs.slide_layouts[5])
	if slide.shapes.title:
		slide.shapes.title.text = "Filtering, merging, final boxes"

	_safe_add_picture(slide, det_dir / "02_boxes_before_merge.png", Inches(0.7), Inches(1.6), width=Inches(6.3))
	_add_caption(slide, "After filtering (before merge)", Inches(0.7), Inches(5.05), Inches(6.3), Inches(0.35))

	_safe_add_picture(slide, det_dir / "03_boxes_final.png", Inches(7.1), Inches(1.6), width=Inches(6.3))
	_add_caption(slide, "Final boxes (merged + sorted)", Inches(7.1), Inches(5.05), Inches(6.3), Inches(0.35))

	# Add heuristics list on bottom
	box = slide.shapes.add_textbox(Inches(0.7), Inches(5.55), Inches(12.7), Inches(3.2))
	frame = box.text_frame
	frame.clear()
	p = frame.paragraphs[0]
	p.text = "Heuristics used (filter_white_box): area, min/max size, aspect ratio, text ratio, brightness (LAB L), contrast (ΔL)"
	p.font.size = Pt(16)

	# --- Slide 8: OCR from crops ---
	slide = prs.slides.add_slide(prs.slide_layouts[1])
	_add_bullets(
		slide,
		title="After detection: crop + OCR the price",
		bullets=[
			"For each detected box: extract crop (saved to debug_dir/crops/box_XX.png)",
			"OCR tries multiple binarizations (Otsu, manual threshold, inverted) and multiple PSM modes",
			"Parse price patterns like d.dd; validate numeric range; apply small OCR error corrections",
		],
		note="Code: src/ocr/pipeline/pricebox_detection.py (extract_price_from_crop)",
	)

	output_path.parent.mkdir(parents=True, exist_ok=True)
	prs.save(str(output_path))


def main() -> None:
	out = WORKSPACE_ROOT / "preprocessing_and_pricebox_detection.pptx"
	build_presentation(out)
	print(f"Wrote: {out}")


if __name__ == "__main__":
	main()
